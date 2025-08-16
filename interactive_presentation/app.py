from flask import Flask, render_template, request, jsonify, send_file
import os
import io
import zipfile
import tempfile
import shutil
import json
from datetime import datetime, date
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Inches
import random
from PIL import Image
from PIL.ExifTags import TAGS
import base64

app = Flask(__name__)

# Global variables to store session data
session_data = {
    'current_step': 1,
    'pptx_data': None,
    'slide_analysis': None,
    'placeholders_config': {},
    'processing_details': [],
    'show_details_needed': False,
    'temp_dir': None
}

def add_detail(message, detail_type="info"):
    """إضافة تفصيل جديد إلى قائمة التفاصيل"""
    session_data['processing_details'].append({
        'message': message,
        'type': detail_type
    })
    
    if detail_type in ['error', 'warning']:
        session_data['show_details_needed'] = True

def clear_details():
    """مسح جميع التفاصيل وإعادة تعيين حالة الإظهار"""
    session_data['processing_details'] = []
    session_data['show_details_needed'] = False

def analyze_slide_placeholders(prs):
    """تحليل جميع placeholders في الشريحة الأولى مع ضبط الإحداثيات"""
    if len(prs.slides) == 0:
        return None
    
    first_slide = prs.slides[0]
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    placeholders = {
        'image_placeholders': [],
        'text_placeholders': [],
        'title_placeholders': [],
        'slide_dimensions': {
            'width': slide_width,
            'height': slide_height,
            'width_inches': slide_width / 914400,
            'height_inches': slide_height / 914400
        }
    }
    
    placeholder_id = 0
    def clamp_percent(val):
        # تأكد أن القيمة بين 0 و 100 دائماً
        return max(0, min(val, 100))
    
    for shape in first_slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            left_percent = clamp_percent((shape.left / slide_width) * 100)
            top_percent = clamp_percent((shape.top / slide_height) * 100)
            width_percent = clamp_percent((shape.width / slide_width) * 100)
            height_percent = clamp_percent((shape.height / slide_height) * 100)
            placeholder_info = {
                'id': placeholder_id,
                'type': placeholder_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'left_percent': left_percent,
                'top_percent': top_percent,
                'width_percent': width_percent,
                'height_percent': height_percent,
                'rotation': getattr(shape, 'rotation', 0)
            }
            if placeholder_type == PP_PLACEHOLDER.PICTURE:
                placeholder_info['current_content'] = "صورة"
                placeholders['image_placeholders'].append(placeholder_info)
            elif placeholder_type == PP_PLACEHOLDER.TITLE:
                placeholder_info['current_content'] = shape.text_frame.text if hasattr(shape, 'text_frame') and shape.text_frame.text else "العنوان"
                placeholders['title_placeholders'].append(placeholder_info)
            else:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    placeholder_info['current_content'] = shape.text_frame.text if shape.text_frame.text else f"نص {placeholder_id + 1}"
                    placeholders['text_placeholders'].append(placeholder_info)
            placeholder_id += 1
    for shape in first_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder:
            left_percent = clamp_percent((shape.left / slide_width) * 100)
            top_percent = clamp_percent((shape.top / slide_height) * 100)
            width_percent = clamp_percent((shape.width / slide_width) * 100)
            height_percent = clamp_percent((shape.height / slide_height) * 100)
            image_info = {
                'id': placeholder_id,
                'type': 'regular_image',
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'left_percent': left_percent,
                'top_percent': top_percent,
                'width_percent': width_percent,
                'height_percent': height_percent,
                'rotation': getattr(shape, 'rotation', 0),
                'current_content': "صورة موجودة"
            }
            placeholders['image_placeholders'].append(image_info)
            placeholder_id += 1
    return placeholders

def get_image_date(image_path):
    """استخراج تاريخ التقاط الصورة من metadata"""
    try:
        with Image.open(image_path) as img:
            exifdata = img.getexif()
            for tag_id in exifdata:
                tag = TAGS.get(tag_id, tag_id)
                data = exifdata.get(tag_id)
                
                if tag in ['DateTime', 'DateTimeOriginal', 'DateTimeDigitized']:
                    try:
                        date_obj = datetime.strptime(str(data), '%Y:%m:%d %H:%M:%S')
                        return date_obj.strftime('%Y-%m-%d')
                    except:
                        continue
        
        timestamp = os.path.getmtime(image_path)
        return datetime.fromtimestamp(timestamp).strftime('%Y-%m-%d')
    except:
        return datetime.now().strftime('%Y-%m-%d')

def apply_configured_placeholders(slide, folder_path, folder_name, slide_analysis, placeholders_config):
    """تطبيق الإعدادات المحددة على الشريحة"""
    
    # الحصول على قائمة الصور في المجلد
    imgs = [f for f in os.listdir(folder_path) 
            if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
    imgs.sort()
    
    # تطبيق إعدادات الصور
    image_config = placeholders_config.get('images', {})
    
    # إنشاء قاموس للصور حسب الترتيب المطلوب
    image_assignments = {}
    for config_key, config in image_config.items():
        if config['use'] and config['order'] and config['order'] <= len(imgs):
            image_path = os.path.join(folder_path, imgs[config['order'] - 1])
            placeholder_info = config['placeholder_info']
            
            # العثور على الشكل المقابل في الشريحة الجديدة
            target_shapes = []
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    target_shapes.append(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder:
                    target_shapes.append(shape)
            
            # مطابقة الشكل بناءً على الموقع
            for shape in target_shapes:
                shape_left_percent = (shape.left / slide_analysis['slide_dimensions']['width']) * 100
                shape_top_percent = (shape.top / slide_analysis['slide_dimensions']['height']) * 100
                
                # تحمل اختلاف بسيط في الموقع
                if (abs(shape_left_percent - placeholder_info['left_percent']) < 5 and 
                    abs(shape_top_percent - placeholder_info['top_percent']) < 5):
                    
                    try:
                        if shape.is_placeholder:
                            with open(image_path, 'rb') as img_file:
                                shape.insert_picture(img_file)
                        else:
                            # استبدال الصورة العادية
                            original_left = shape.left
                            original_top = shape.top
                            original_width = shape.width
                            original_height = shape.height
                            
                            shape_element = shape._element
                            shape_element.getparent().remove(shape_element)
                            
                            slide.shapes.add_picture(image_path, original_left, original_top, original_width, original_height)
                        
                        add_detail(f"✅ تم استبدال الصورة {config['order']}: {os.path.basename(image_path)}", "success")
                        break
                    except Exception as e:
                        add_detail(f"❌ فشل في استبدال الصورة: {e}", "error")
    
    # تطبيق إعدادات النصوص
    text_config = placeholders_config.get('texts', {})
    
    text_shapes = []
    for shape in slide.shapes:
        if (shape.is_placeholder and 
            shape.placeholder_format.type not in [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.TITLE] and
            hasattr(shape, 'text_frame') and shape.text_frame):
            text_shapes.append(shape)
    
    text_index = 0
    for config_key, config in text_config.items():
        if text_index < len(text_shapes):
            shape = text_shapes[text_index]
            
            try:
                if config['type'] == "ترك فارغ":
                    shape.text_frame.text = ""
                    
                elif config['type'] == "نص ثابت":
                    if config['value']:
                        shape.text_frame.text = config['value']
                        
                elif config['type'] == "تاريخ":
                    if config['value'] == "today":
                        date_text = datetime.now().strftime('%Y-%m-%d')
                    else:
                        date_text = config['value']
                    shape.text_frame.text = date_text
                    
                elif config['type'] == "تاريخ الصورة" and imgs:
                    first_image_path = os.path.join(folder_path, imgs[0])
                    image_date = get_image_date(first_image_path)
                    shape.text_frame.text = image_date
                    
                elif config['type'] == "اسم المجلد":
                    shape.text_frame.text = folder_name
                
                add_detail(f"✅ تم تطبيق النص: {config['type']}", "success")
                
            except Exception as e:
                add_detail(f"⚠ خطأ في تطبيق النص: {e}", "warning")
            
            text_index += 1
    
    # تطبيق العنوان (اسم المجلد)
    title_shapes = [
        shape for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
    ]
    
    if title_shapes:
        title_shapes[0].text = folder_name
        add_detail(f"✅ تم تحديث العنوان: {folder_name}", "success")

@app.route('/')
def index():
    # Reset session data when accessing the home page
    global session_data
    session_data = {
        'current_step': 1,
        'pptx_data': None,
        'slide_analysis': None,
        'placeholders_config': {},
        'processing_details': [],
        'show_details_needed': False,
        'temp_dir': None
    }
    return render_template('index.html')

@app.route('/upload-pptx', methods=['POST'])
def upload_pptx():
    if 'pptx_file' not in request.files:
        return jsonify({'success': False, 'error': 'No file uploaded'})
    
    pptx_file = request.files['pptx_file']
    
    if pptx_file.filename == '':
        return jsonify({'success': False, 'error': 'No file selected'})
    
    if not pptx_file.filename.endswith('.pptx'):
        return jsonify({'success': False, 'error': 'File must be a .pptx file'})
    
    try:
        # Save the file data
        pptx_data = pptx_file.read()
        session_data['pptx_data'] = pptx_data
        
        # Analyze the slide
        prs = Presentation(io.BytesIO(pptx_data))
        slide_analysis = analyze_slide_placeholders(prs)
        
        if slide_analysis:
            session_data['slide_analysis'] = slide_analysis
            session_data['current_step'] = 2
            return jsonify({
                'success': True, 
                'slide_analysis': slide_analysis,
                'redirect': '/configure'
            })
        else:
            return jsonify({'success': False, 'error': 'No slides found in the file or analysis error'})
            
    except Exception as e:
        return jsonify({'success': False, 'error': f'Error analyzing file: {str(e)}'})

@app.route('/configure')
def configure():
    if session_data['current_step'] < 2 or not session_data['slide_analysis']:
        return render_template('error.html', message='Please upload a PowerPoint file first')
    
    return render_template('configure.html', 
                          slide_analysis=session_data['slide_analysis'],
                          current_step=session_data['current_step'])

@app.route('/save-config', methods=['POST'])
def save_config():
    config_data = request.json
    session_data['placeholders_config'] = config_data
    session_data['current_step'] = 3
    return jsonify({'success': True, 'redirect': '/process'})

@app.route('/use-previous-settings', methods=['POST'])
def use_previous_settings():
    try:
        config_data = request.json
        if not config_data:
            return jsonify({'success': False, 'error': 'No configuration data provided'})
        
        # Store the configuration in the session
        session_data['placeholders_config'] = config_data
        session_data['current_step'] = 3
        session_data['using_previous_settings'] = True
        
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/process')
def process():
    if session_data['current_step'] < 3 or not session_data['placeholders_config']:
        return render_template('error.html', message='Please configure placeholders first')
    
    return render_template('process.html', 
                          placeholders_config=session_data['placeholders_config'],
                          current_step=session_data['current_step'])

@app.route('/upload-zip', methods=['POST'])
def upload_zip():
    if 'zip_file' not in request.files:
        return jsonify({'success': False, 'error': 'No file uploaded'})
    
    zip_file = request.files['zip_file']
    
    if zip_file.filename == '':
        return jsonify({'success': False, 'error': 'No file selected'})
    
    if not zip_file.filename.endswith('.zip'):
        return jsonify({'success': False, 'error': 'File must be a .zip file'})
    
    options = {
        'image_order': request.form.get('image_order', 'alphabetical'),
        'skip_empty_folders': request.form.get('skip_empty_folders', 'true') == 'true'
    }
    
    clear_details()
    
    try:
        # Extract the zip file
        zip_bytes = io.BytesIO(zip_file.read())
        
        # Clean up previous temp directory if exists
        if session_data['temp_dir'] and os.path.exists(session_data['temp_dir']):
            shutil.rmtree(session_data['temp_dir'])
        
        # Create new temp directory
        temp_dir = tempfile.mkdtemp()
        session_data['temp_dir'] = temp_dir
        
        with zipfile.ZipFile(zip_bytes, "r") as zip_ref:
            zip_ref.extractall(temp_dir)
        
        add_detail("📂 تم استخراج الملف المضغوط بنجاح", "success")
        
        # Find folders with images
        all_items = os.listdir(temp_dir)
        folder_paths = []
        
        for item in all_items:
            item_path = os.path.join(temp_dir, item)
            if os.path.isdir(item_path):
                imgs_in_folder = [f for f in os.listdir(item_path) 
                                if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                if imgs_in_folder:
                    folder_paths.append(item_path)
                    add_detail(f"📁 المجلد '{item}' يحتوي على {len(imgs_in_folder)} صورة", "info")
                elif not options['skip_empty_folders']:
                    add_detail(f"⚠ المجلد '{item}' فارغ من الصور", "warning")
        
        if not folder_paths:
            return jsonify({
                'success': False, 
                'error': 'لا توجد مجلدات تحتوي على صور في الملف المضغوط.',
                'details': session_data['processing_details']
            })
        
        folder_paths.sort()
        add_detail(f"✅ تم العثور على {len(folder_paths)} مجلد يحتوي على صور", "success")
        
        # Load PowerPoint file
        prs = Presentation(io.BytesIO(session_data['pptx_data']))
        
        if len(prs.slides) == 0:
            return jsonify({
                'success': False, 
                'error': 'لا توجد شرائح في ملف PowerPoint',
                'details': session_data['processing_details']
            })
        
        first_slide = prs.slides[0]
        slide_layout = first_slide.slide_layout
        
        # Process slides
        total_processed = 0
        created_slides = 0
        
        for folder_idx, folder_path in enumerate(folder_paths):
            folder_name = os.path.basename(folder_path)
            
            try:
                # Sort images in the folder
                imgs = [f for f in os.listdir(folder_path) 
                       if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                
                if options['image_order'] == 'random':
                    random.shuffle(imgs)
                    add_detail(f"🔀 تم ترتيب صور المجلد {folder_name} عشوائياً", "info")
                else:
                    imgs.sort()
                    add_detail(f"📋 تم ترتيب صور المجلد {folder_name} أبجدياً", "info")
                
                # Create new slide
                new_slide = prs.slides.add_slide(slide_layout)
                created_slides += 1
                
                # Apply configured placeholders
                apply_configured_placeholders(
                    new_slide, 
                    folder_path, 
                    folder_name, 
                    session_data['slide_analysis'],
                    session_data['placeholders_config']
                )
                
                total_processed += len(imgs)
                add_detail(f"✅ تم إنشاء شريحة للمجلد '{folder_name}' مع {len(imgs)} صورة", "success")
            
            except Exception as e:
                add_detail(f"❌ خطأ في معالجة المجلد {folder_name}: {str(e)}", "error")
        
        # Save the file
        output_filename = f"PowerPoint_Updated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = os.path.join(temp_dir, output_filename)
        prs.save(output_path)
        
        return jsonify({
            'success': True,
            'message': 'تم الانتهاء من المعالجة بنجاح!',
            'details': session_data['processing_details'],
            'stats': {
                'created_slides': created_slides,
                'processed_folders': len(folder_paths),
                'total_images': total_processed
            },
            'output_filename': output_filename
        })
        
    except Exception as e:
        add_detail(f"❌ خطأ عام أثناء المعالجة: {str(e)}", "error")
        return jsonify({
            'success': False,
            'error': f'Error processing files: {str(e)}',
            'details': session_data['processing_details']
        })

@app.route('/download/<filename>')
def download_file(filename):
    if not session_data['temp_dir'] or not os.path.exists(session_data['temp_dir']):
        return render_template('error.html', message='No processed file available')
    
    file_path = os.path.join(session_data['temp_dir'], filename)
    
    if not os.path.exists(file_path):
        return render_template('error.html', message='File not found')
    
    return send_file(file_path, as_attachment=True)

@app.route('/reset')
def reset():
    global session_data
    # Clean up temp directory
    if session_data['temp_dir'] and os.path.exists(session_data['temp_dir']):
        try:
            shutil.rmtree(session_data['temp_dir'])
        except:
            pass
    
    # Reset session data
    session_data = {
        'current_step': 1,
        'pptx_data': None,
        'slide_analysis': None,
        'placeholders_config': {},
        'processing_details': [],
        'show_details_needed': False,
        'temp_dir': None
    }
    
    return jsonify({'success': True, 'redirect': '/'})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)